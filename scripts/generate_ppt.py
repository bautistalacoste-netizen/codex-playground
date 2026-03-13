from __future__ import annotations

from dataclasses import dataclass, field
from io import BytesIO
from pathlib import Path
import re
from typing import List, Sequence

import fitz  # PyMuPDF
from PIL import Image
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

INPUT_PATH = Path("content/model_walkthrough.pdf")
OUTPUT_PATH = Path("output/model_walkthrough_presentation.pptx")


@dataclass
class ExtractedImage:
    data: bytes
    ext: str
    width: int
    height: int
    is_chart_like: bool = False


@dataclass
class PageContent:
    page_number: int
    title: str
    bullets: List[str] = field(default_factory=list)
    images: List[ExtractedImage] = field(default_factory=list)


def normalize_line(line: str) -> str:
    cleaned = re.sub(r"\s+", " ", line).strip()
    return cleaned.strip("•-") if cleaned else ""


def detect_chart_like_image(width: int, height: int) -> bool:
    """Heuristic to tag visual content that may include charts/diagrams."""
    if width <= 0 or height <= 0:
        return False

    area = width * height
    aspect = width / height
    return area > 60_000 and 0.75 <= aspect <= 2.8


def extract_page_text(page: fitz.Page) -> tuple[str, List[str]]:
    text = page.get_text("text")
    lines = [normalize_line(line) for line in text.splitlines()]
    lines = [line for line in lines if line]

    if not lines:
        return f"Page {page.number + 1}", []

    title = lines[0][:120]
    bullet_lines: List[str] = []

    for raw in lines[1:]:
        if len(raw) < 3:
            continue

        # Keep naturally-bulleted lines and sentences.
        if re.match(r"^(\d+[.)]|[a-zA-Z][.)])\s+", raw):
            bullet_lines.append(re.sub(r"^(\d+[.)]|[a-zA-Z][.)])\s+", "", raw))
        else:
            bullet_lines.append(raw)

    # Keep slide text concise.
    return title, bullet_lines[:8]


def extract_page_images(doc: fitz.Document, page: fitz.Page) -> List[ExtractedImage]:
    extracted: List[ExtractedImage] = []
    image_refs = page.get_images(full=True)

    for image_ref in image_refs:
        xref = image_ref[0]
        image_info = doc.extract_image(xref)
        image_bytes = image_info.get("image")
        if not image_bytes:
            continue

        ext = image_info.get("ext", "png")
        with Image.open(BytesIO(image_bytes)) as img:
            width, height = img.size

        extracted.append(
            ExtractedImage(
                data=image_bytes,
                ext=ext,
                width=width,
                height=height,
                is_chart_like=detect_chart_like_image(width, height),
            )
        )

    return extracted


def extract_pdf_content(pdf_path: Path) -> List[PageContent]:
    pages: List[PageContent] = []

    with fitz.open(pdf_path) as doc:
        for page in doc:
            title, bullets = extract_page_text(page)
            images = extract_page_images(doc, page)

            pages.append(
                PageContent(
                    page_number=page.number + 1,
                    title=title,
                    bullets=bullets,
                    images=images,
                )
            )

    return pages


def apply_corporate_style(prs: Presentation) -> None:
    prs.slide_width = Inches(13.33)
    prs.slide_height = Inches(7.5)


def add_title_slide(prs: Presentation, title: str, subtitle: str) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = title
    slide.placeholders[1].text = subtitle

    title_p = slide.shapes.title.text_frame.paragraphs[0]
    title_p.font.name = "Calibri"
    title_p.font.size = Pt(40)
    title_p.font.bold = True
    title_p.font.color.rgb = RGBColor(15, 56, 107)

    subtitle_p = slide.placeholders[1].text_frame.paragraphs[0]
    subtitle_p.font.name = "Calibri"
    subtitle_p.font.size = Pt(20)
    subtitle_p.font.color.rgb = RGBColor(90, 90, 90)


def add_agenda_slide(prs: Presentation, pages: Sequence[PageContent]) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "Agenda"
    text_frame = slide.shapes.placeholders[1].text_frame
    text_frame.clear()

    for idx, page in enumerate(pages[:12], start=1):
        paragraph = text_frame.paragraphs[0] if idx == 1 else text_frame.add_paragraph()
        paragraph.text = f"{idx}. {page.title}"
        paragraph.font.name = "Calibri"
        paragraph.font.size = Pt(20)


def add_page_slide(prs: Presentation, page: PageContent) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[5])

    title_shape = slide.shapes.title
    title_shape.text = page.title
    title_p = title_shape.text_frame.paragraphs[0]
    title_p.font.name = "Calibri"
    title_p.font.bold = True
    title_p.font.size = Pt(30)
    title_p.font.color.rgb = RGBColor(15, 56, 107)

    text_box = slide.shapes.add_textbox(Inches(0.7), Inches(1.45), Inches(7.0), Inches(5.6))
    tf = text_box.text_frame
    tf.word_wrap = True

    bullets = page.bullets if page.bullets else [f"Extracted from PDF page {page.page_number}."]
    for idx, bullet in enumerate(bullets[:10]):
        para = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
        para.text = bullet
        para.level = 0
        para.font.name = "Calibri"
        para.font.size = Pt(18)
        para.font.color.rgb = RGBColor(50, 50, 50)

    prioritized_images = sorted(page.images, key=lambda img: (not img.is_chart_like, -(img.width * img.height)))

    if prioritized_images:
        image = prioritized_images[0]
        image_stream = BytesIO(image.data)
        left = Inches(8.0)
        top = Inches(1.55)
        max_w = 4.9
        max_h = 4.9

        ratio = image.width / image.height if image.height else 1
        if ratio >= 1:
            width = Inches(max_w)
            height = Inches(max_w / ratio)
            if height > Inches(max_h):
                height = Inches(max_h)
                width = Inches(max_h * ratio)
        else:
            height = Inches(max_h)
            width = Inches(max_h * ratio)
            if width > Inches(max_w):
                width = Inches(max_w)
                height = Inches(max_w / ratio)

        slide.shapes.add_picture(image_stream, left, top, width=width, height=height)

        if image.is_chart_like:
            label_box = slide.shapes.add_textbox(Inches(8.0), Inches(6.6), Inches(4.9), Inches(0.4))
            label_frame = label_box.text_frame
            label_frame.text = "Extracted chart/visual"
            label_p = label_frame.paragraphs[0]
            label_p.alignment = PP_ALIGN.RIGHT
            label_p.font.name = "Calibri"
            label_p.font.size = Pt(12)
            label_p.font.color.rgb = RGBColor(100, 100, 100)


def add_closing_slide(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    slide.shapes.title.text = "Next Steps"

    text_box = slide.shapes.add_textbox(Inches(1.0), Inches(2.1), Inches(11.0), Inches(2.5))
    tf = text_box.text_frame
    tf.text = "Review extracted content, validate chart placement, and refine speaker notes."
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    p.font.name = "Calibri"
    p.font.size = Pt(28)
    p.font.color.rgb = RGBColor(15, 56, 107)


def resolve_input_pdf(default_path: Path) -> Path:
    if default_path.exists():
        return default_path

    candidates = sorted(Path("content").glob("*.pdf"))
    if len(candidates) == 1:
        return candidates[0]

    raise FileNotFoundError(
        "Missing input file: content/model_walkthrough.pdf. "
        "Place the walkthrough PDF there (or keep only one PDF inside content/)."
    )


def generate_presentation(input_path: Path = INPUT_PATH, output_path: Path = OUTPUT_PATH) -> Path:
    pdf_path = resolve_input_pdf(input_path)
    pages = extract_pdf_content(pdf_path)
    if not pages:
        raise ValueError(f"No readable content found in PDF: {pdf_path}")

    prs = Presentation()
    apply_corporate_style(prs)
    add_title_slide(prs, "Model Walkthrough", f"Generated from {pdf_path.name}")
    add_agenda_slide(prs, pages)

    for page in pages:
        add_page_slide(prs, page)

    add_closing_slide(prs)

    output_path.parent.mkdir(parents=True, exist_ok=True)
    prs.save(output_path)
    return output_path


def main() -> None:
    generated_path = generate_presentation()
    print(f"Presentation generated at: {generated_path}")


if __name__ == "__main__":
    main()
